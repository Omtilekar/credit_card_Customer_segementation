"""
Professional academic PowerPoint — Credit Card Customer Segmentation
University of Delaware · Introduction to Data Mining
20 slides, clean navy/white academic theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

# ── Output / asset paths ──────────────────────────────────────────────────────
PLOTS = "outputs/plots"
OUT   = "outputs/results/CC_Segmentation_Presentation.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x3A, 0x5C)   # primary brand navy
STEEL     = RGBColor(0x2E, 0x86, 0xAB)   # steel blue accent
LIGHT_BG  = RGBColor(0xEB, 0xF5, 0xFB)   # very light blue content bg
PALE      = RGBColor(0xD6, 0xEA, 0xF8)   # pale blue for rows/cards
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEXT      = RGBColor(0x1C, 0x28, 0x33)   # near-black body text
LGRAY     = RGBColor(0x7F, 0x8C, 0x8D)   # secondary text
DGRAY     = RGBColor(0x44, 0x44, 0x44)
GREEN     = RGBColor(0x1E, 0x84, 0x49)
ORANGE    = RGBColor(0xD3, 0x54, 0x00)
RED       = RGBColor(0xC0, 0x39, 0x2B)
YELLOW_BG = RGBColor(0xFF, 0xF9, 0xE6)
TEAL      = RGBColor(0x17, 0xA5, 0x89)

# Cluster-specific accent colours
CC = [STEEL, RED, GREEN, ORANGE]          # C0 C1 C2 C3

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)
BLANK   = None   # set after Presentation() below

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]   # completely blank

# ── Low-level drawing helpers ─────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=0.5,
         rounded=False, radius_pct=8):
    """Add a filled rectangle (or rounded rectangle)."""
    shape_type = 5 if rounded else 1   # 5 = ROUNDED_RECTANGLE, 1 = RECTANGLE
    shp = slide.shapes.add_shape(shape_type,
                                  Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    if rounded:
        sp_pr = shp.element.find(qn('p:spPr'))
        if sp_pr is None:
            sp_pr = shp.element.spPr
        pst = sp_pr.find(qn('a:prstGeom'))
        if pst is not None:
            av_lst = pst.find(qn('a:avLst'))
            if av_lst is not None:
                gd = av_lst.find(qn('a:gd'))
                if gd is not None:
                    gd.set('fmla', f'val {radius_pct * 1000}')
    return shp


def txt(slide, text, l, t, w, h, size=14, bold=False, italic=False,
        color=TEXT, align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    """Add a textbox."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def img(slide, path, l, t, w):
    """Add an image scaled proportionally to width w (inches)."""
    if not os.path.exists(path):
        return
    slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))


def bullet_frame(slide, items, l, t, w, h, size=13, color=TEXT,
                 bold_first=False, indent="  "):
    """Multi-line bullet textbox where each item is one paragraph."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{indent}•  {item}"
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold_first and i == 0
        run.font.color.rgb = color
    return tb


def slide_header(slide, title, sub=None, show_num=None):
    """Standard navy header bar at top of slide."""
    rect(slide, 0, 0, 13.33, 0.85, fill=NAVY)
    # thin steel-blue accent line
    rect(slide, 0, 0.85, 13.33, 0.04, fill=STEEL)
    txt(slide, title, 0.3, 0.07, 11.5, 0.72,
        size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if sub:
        txt(slide, sub, 0.3, 0.62, 10.0, 0.26,
            size=11, color=RGBColor(0xA9, 0xCC, 0xE3), align=PP_ALIGN.LEFT)
    if show_num:
        txt(slide, str(show_num), 12.7, 0.12, 0.5, 0.5,
            size=11, color=RGBColor(0xA9, 0xCC, 0xE3), align=PP_ALIGN.RIGHT)
    return slide


def section_label(slide, text, l=0.3, t=0.95):
    """Small teal tag label under the header bar."""
    rect(slide, l, t, len(text)*0.115 + 0.25, 0.3, fill=STEEL, rounded=True)
    txt(slide, text, l + 0.12, t + 0.03, len(text)*0.115 + 0.12, 0.25,
        size=9.5, bold=True, color=WHITE)


def info_card(slide, title, body_lines, l, t, w, h,
              accent=NAVY, bg=LIGHT_BG, title_size=12, body_size=11):
    """A titled card box with bullet lines inside."""
    rect(slide, l, t, w, h, fill=bg, line_color=accent, line_w=1.0)
    rect(slide, l, t, w, 0.38, fill=accent)
    txt(slide, title, l + 0.12, t + 0.05, w - 0.2, 0.3,
        size=title_size, bold=True, color=WHITE)
    by = t + 0.44
    line_h = (h - 0.52) / max(len(body_lines), 1)
    for line in body_lines:
        txt(slide, f"• {line}", l + 0.14, by, w - 0.22, line_h + 0.05,
            size=body_size, color=TEXT)
        by += line_h


def metric_card(slide, label, value, l, t, w=2.2, h=1.0,
                accent=STEEL, value_size=28):
    """Compact KPI card."""
    rect(slide, l, t, w, h, fill=WHITE, line_color=accent, line_w=1.2)
    rect(slide, l, t, w, 0.28, fill=accent)
    txt(slide, label, l + 0.08, t + 0.03, w - 0.12, 0.24,
        size=8.5, bold=True, color=WHITE)
    txt(slide, value, l + 0.05, t + 0.30, w - 0.1, h - 0.38,
        size=value_size, bold=True, color=accent, align=PP_ALIGN.CENTER)


def slide_number(slide, n):
    txt(slide, f"Slide {n}", 12.3, 7.15, 0.95, 0.28,
        size=8, color=LGRAY, align=PP_ALIGN.RIGHT)


def flow_step(slide, label, x, y, w=1.9, h=0.72, fill=NAVY, tsize=10):
    rect(slide, x, y, w, h, fill=fill, rounded=True)
    txt(slide, label, x + 0.06, y + 0.1, w - 0.1, h - 0.14,
        size=tsize, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def arrow_right(slide, x, y):
    """Draw a right-pointing arrow connector."""
    conn = slide.shapes.add_connector(1, Inches(x), Inches(y),
                                       Inches(x + 0.35), Inches(y))
    conn.line.color.rgb = STEEL
    conn.line.width = Pt(2)


def arrow_down(slide, x, y, dy=0.3):
    conn = slide.shapes.add_connector(1, Inches(x), Inches(y),
                                       Inches(x), Inches(y + dy))
    conn.line.color.rgb = STEEL
    conn.line.width = Pt(2)


def pptx_table(slide, headers, rows, l, t, w, h,
               col_ratios=None, hdr_color=NAVY):
    """Add a formatted table."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl = slide.shapes.add_table(n_rows, n_cols,
                                  Inches(l), Inches(t),
                                  Inches(w), Inches(h)).table

    # column widths
    if col_ratios:
        total = sum(col_ratios)
        for i, r in enumerate(col_ratios):
            tbl.columns[i].width = Inches(w * r / total)

    def cell_fmt(cell, text, bold=False, bg=None, color=TEXT,
                 size=9, align=PP_ALIGN.CENTER):
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = color
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # header row
    for j, h_text in enumerate(headers):
        c = tbl.cell(0, j)
        cell_fmt(c, h_text, bold=True, bg=hdr_color, color=WHITE, size=9)

    # data rows
    for i, row in enumerate(rows):
        bg = PALE if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            is_left = (j == 0)
            cell_fmt(c, str(val), bg=bg, size=9,
                     align=PP_ALIGN.LEFT if is_left else PP_ALIGN.CENTER)
    return tbl


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
def slide_01():
    sl = prs.slides.add_slide(BLANK)
    # Full background
    rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    # Diagonal accent block (right side)
    rect(sl, 9.2, 0, 4.13, 7.5, fill=RGBColor(0x1E, 0x44, 0x6A))
    rect(sl, 9.1, 0, 0.12, 7.5, fill=STEEL)

    # University wordmark area
    rect(sl, 0.5, 0.45, 3.5, 0.5, fill=RGBColor(0x25, 0x4E, 0x78), rounded=True)
    txt(sl, "University of Delaware", 0.62, 0.50, 3.26, 0.4,
        size=11, bold=True, color=RGBColor(0xA9, 0xCC, 0xE3))

    # Main title
    txt(sl, "Credit Card Customer", 0.5, 1.3, 8.5, 1.1,
        size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(sl, "Segmentation", 0.5, 2.35, 8.5, 0.9,
        size=42, bold=True, color=RGBColor(0x5D, 0xAD, 0xE8), align=PP_ALIGN.LEFT)
    txt(sl, "Using Clustering Techniques", 0.5, 3.2, 8.5, 0.7,
        size=22, color=RGBColor(0xB2, 0xCC, 0xE2), align=PP_ALIGN.LEFT)

    # Decorative dots
    for i, y in enumerate([1.5, 2.5, 3.5, 4.5]):
        dot = sl.shapes.add_shape(9, Inches(0.22), Inches(y), Inches(0.1), Inches(0.1))
        dot.fill.solid()
        dot.fill.fore_color.rgb = STEEL
        dot.line.fill.background()

    # Info panel on right
    rect(sl, 9.4, 1.5, 3.65, 4.8, fill=RGBColor(0x16, 0x34, 0x54))
    txt(sl, "Course", 9.6, 1.65, 3.3, 0.3, size=10, color=STEEL, bold=True)
    txt(sl, "Introduction to Data Mining", 9.6, 1.9, 3.3, 0.5,
        size=13, bold=True, color=WHITE)
    rect(sl, 9.4, 2.42, 3.65, 0.02, fill=STEEL)
    txt(sl, "Team", 9.6, 2.55, 3.3, 0.3, size=10, color=STEEL, bold=True)
    for i, m in enumerate(["Team Member 1", "Team Member 2",
                            "Team Member 3", "Team Member 4"]):
        txt(sl, f"  {m}", 9.6, 2.82 + i * 0.38, 3.3, 0.36, size=12, color=WHITE)
    rect(sl, 9.4, 4.4, 3.65, 0.02, fill=STEEL)
    txt(sl, "Dataset", 9.6, 4.52, 3.3, 0.28, size=10, color=STEEL, bold=True)
    txt(sl, "CC_GENERAL.csv", 9.6, 4.78, 3.3, 0.35, size=12, color=WHITE)
    txt(sl, "8,950 customers  ·  17 features", 9.6, 5.1, 3.3, 0.35,
        size=10.5, color=RGBColor(0xA9, 0xCC, 0xE3))

    # Bottom bar
    rect(sl, 0, 6.9, 13.33, 0.6, fill=RGBColor(0x0F, 0x24, 0x38))
    txt(sl, "University of Delaware  ·  Introduction to Data Mining  ·  May 2026",
        0.4, 6.97, 12.53, 0.4, size=10,
        color=RGBColor(0x6C, 0x93, 0xB5), align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement / Objective
# ═════════════════════════════════════════════════════════════════════════════
def slide_02():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Problem Statement & Objective", "Why segmentation matters — and what we aim to do", 2)
    section_label(sl, "MOTIVATION")
    slide_number(sl, 2)

    # Problem box
    rect(sl, 0.3, 1.35, 6.1, 5.6, fill=LIGHT_BG, line_color=NAVY, line_w=1.0)
    rect(sl, 0.3, 1.35, 6.1, 0.42, fill=NAVY)
    txt(sl, "The Problem", 0.5, 1.38, 5.7, 0.36, size=14, bold=True, color=WHITE)
    problems = [
        "Banks serve millions of customers with very different needs.",
        "A single product strategy fails to address diverse financial behaviors.",
        "Without segmentation, marketing spend is wasted on the wrong audience.",
        "Credit risk cannot be managed effectively without knowing customer types.",
        "Low-engagement customers churn silently with no targeted intervention.",
        "High-value customers receive the same treatment as low-activity ones.",
    ]
    by = 1.88
    for pb in problems:
        rect(sl, 0.42, by, 5.86, 0.64, fill=WHITE, line_color=PALE, line_w=0.5)
        txt(sl, "✗", 0.52, by + 0.14, 0.28, 0.36, size=12, bold=True,
            color=RED)
        txt(sl, pb, 0.84, by + 0.07, 5.3, 0.5, size=10.5, color=TEXT)
        by += 0.73

    # Objective box
    rect(sl, 6.73, 1.35, 6.3, 2.6, fill=LIGHT_BG, line_color=STEEL, line_w=1.0)
    rect(sl, 6.73, 1.35, 6.3, 0.42, fill=STEEL)
    txt(sl, "Our Objective", 6.93, 1.38, 5.9, 0.36, size=14, bold=True, color=WHITE)
    txt(sl,
        "Apply unsupervised machine learning to group credit card "
        "customers into meaningful behavioral segments based on their "
        "spending, payment, and borrowing patterns — without any "
        "predefined labels.",
        6.93, 1.87, 5.9, 1.95, size=11.5, color=TEXT)

    # Approach box
    rect(sl, 6.73, 4.1, 6.3, 2.85, fill=LIGHT_BG, line_color=TEAL, line_w=1.0)
    rect(sl, 6.73, 4.1, 6.3, 0.42, fill=TEAL)
    txt(sl, "Our Approach", 6.93, 4.13, 5.9, 0.36, size=14, bold=True, color=WHITE)
    steps = ["Explore and clean the CC_GENERAL dataset",
             "Apply K-Means, DBSCAN, GMM & Agglomerative",
             "Compare algorithms using 3 internal metrics",
             "Interpret and name the final customer segments"]
    ay = 4.62
    for i, s in enumerate(steps):
        rect(sl, 6.88, ay, 0.34, 0.34, fill=TEAL, rounded=True)
        txt(sl, str(i + 1), 6.89, ay + 0.04, 0.32, 0.28, size=10,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, s, 7.3, ay + 0.03, 5.55, 0.36, size=11, color=TEXT)
        ay += 0.52


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Real-World Applications
# ═════════════════════════════════════════════════════════════════════════════
def slide_03():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Real-World Applications", "How customer segmentation creates business value", 3)
    section_label(sl, "IMPACT")
    slide_number(sl, 3)

    apps = [
        (NAVY, "Marketing\nPersonalization",
         ["Tailor messaging to each segment's spending habits",
          "Send credit limit upgrade offers to Premium Transactors",
          "Re-engagement emails to Dormant customers"]),
        (STEEL, "Credit Risk\nMonitoring",
         ["Cash-Advance Revolvers carry the highest default risk",
          "Set differentiated credit limits per segment",
          "Flag high-balance revolvers for early intervention"]),
        (TEAL, "Customer\nRetention",
         ["Identify Dormant customers before they close accounts",
          "Offer loyalty perks to Premium Transactors",
          "Prevent churn with segment-specific incentives"]),
        (GREEN, "Product\nRecommendation",
         ["Revolvers → balance-transfer & consolidation loans",
          "Transactors → rewards & cashback credit cards",
          "Instalment Buyers → BNPL & low-interest plans"]),
        (ORANGE, "Reward Program\nTargeting",
         ["Match reward type to spending behaviour",
          "Travel miles for high-spend transactors",
          "Instalment discounts for frequent instalment users"]),
    ]
    xs = [0.3, 2.9, 5.5, 8.1, 10.7]
    for i, (color, title, bullets) in enumerate(apps):
        x = xs[i]
        rect(sl, x, 1.35, 2.35, 5.6, fill=WHITE, line_color=color, line_w=1.2)
        rect(sl, x, 1.35, 2.35, 0.7, fill=color)
        txt(sl, title, x + 0.1, 1.38, 2.15, 0.62,
            size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        by = 2.18
        for b in bullets:
            rect(sl, x + 0.1, by, 0.22, 0.22, fill=color, rounded=True)
            txt(sl, "→", x + 0.12, by + 0.02, 0.2, 0.2, size=9,
                bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            txt(sl, b, x + 0.38, by, 1.85, 0.55, size=9.5, color=TEXT)
            by += 0.75


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Data Information
# ═════════════════════════════════════════════════════════════════════════════
def slide_04():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Data Information", "CC_GENERAL.csv — Kaggle Credit Card Dataset", 4)
    section_label(sl, "DATASET")
    slide_number(sl, 4)

    # KPI row
    kpis = [("Customers", "8,950"), ("Raw Features", "18"),
            ("After Cleaning", "17"), ("Missing Values", "314")]
    kpi_colors = [NAVY, STEEL, TEAL, ORANGE]
    for i, ((lbl, val), col) in enumerate(zip(kpis, kpi_colors)):
        metric_card(sl, lbl, val, 0.3 + i * 3.25, 1.35, w=3.0, h=1.1,
                    accent=col, value_size=30)

    # Feature table
    txt(sl, "Key Features", 0.3, 2.62, 5.0, 0.35, size=13, bold=True, color=NAVY)
    headers = ["Feature", "Type", "Description"]
    rows = [
        ["BALANCE",             "Float", "Avg monthly outstanding balance ($)"],
        ["PURCHASES",           "Float", "Total purchases made in period ($)"],
        ["ONEOFF_PURCHASES",    "Float", "One-off (single) purchase amount ($)"],
        ["INSTALLMENTS_PURCH.", "Float", "Instalment purchase amount ($)"],
        ["CASH_ADVANCE",        "Float", "Total cash advance amount ($)"],
        ["CREDIT_LIMIT",        "Float", "Customer credit limit ($)"],
        ["PAYMENTS",            "Float", "Total payments made ($)"],
        ["MINIMUM_PAYMENTS",    "Float", "Minimum payments required ($)"],
        ["PRC_FULL_PAYMENT",    "Float [0,1]", "Fraction of months paid in full"],
        ["PURCHASES_FREQUENCY", "Float [0,1]", "Fraction of months with purchases"],
        ["CASH_ADV_FREQUENCY",  "Float [0,1]", "Fraction of months with advances"],
        ["TENURE",              "Integer", "Months as a customer"],
    ]
    pptx_table(sl, headers, rows, 0.3, 3.05, 8.5, 4.2,
               col_ratios=[2.2, 1.5, 4.0])

    # Source & notes panel
    rect(sl, 9.0, 2.62, 4.0, 4.63, fill=LIGHT_BG, line_color=NAVY, line_w=0.8)
    rect(sl, 9.0, 2.62, 4.0, 0.4, fill=NAVY)
    txt(sl, "Dataset Notes", 9.15, 2.65, 3.7, 0.34, size=12, bold=True, color=WHITE)
    notes = ["Source: Kaggle (CC_GENERAL Dataset)",
             "One row per customer (6-month window)",
             "CUST_ID dropped — not a feature",
             "Missing: MINIMUM_PAYMENTS (313 rows)",
             "Missing: CREDIT_LIMIT (1 row)",
             "→ Both filled with column median",
             "10 monetary features are right-skewed",
             "→ log₁p transform applied before scaling",
             "6 frequency features already in [0, 1]",
             "→ Kept on original scale"]
    ny = 3.12
    for n in notes:
        color = ORANGE if n.startswith("Missing") else (TEAL if n.startswith("→") else TEXT)
        txt(sl, f"  {n}", 9.1, ny, 3.75, 0.37, size=10, color=color)
        ny += 0.39


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — How the Data Looks
# ═════════════════════════════════════════════════════════════════════════════
def slide_05():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "How the Data Looks", "Feature distributions, missing values, and correlations", 5)
    section_label(sl, "EXPLORATORY DATA ANALYSIS")
    slide_number(sl, 5)

    txt(sl, "Feature Distributions", 0.3, 1.35, 7.5, 0.34, size=12, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/feature_distributions.png", 0.3, 1.72, 8.4)

    # Right panel
    txt(sl, "Correlation Heatmap", 8.95, 1.35, 4.1, 0.34, size=12, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/correlation_heatmap.png", 8.95, 1.72, 4.1)

    # Missing value summary strip
    rect(sl, 0.3, 6.55, 12.73, 0.75, fill=LIGHT_BG, line_color=ORANGE, line_w=0.8)
    txt(sl, "Missing Values:", 0.5, 6.63, 2.5, 0.4, size=11, bold=True, color=ORANGE)
    txt(sl, "MINIMUM_PAYMENTS: 313 rows (3.5%) — imputed with median  |  "
            "CREDIT_LIMIT: 1 row (<0.01%) — imputed with median  |  "
            "All other features: complete  |  No rows removed",
        3.0, 6.63, 9.9, 0.5, size=10.5, color=TEXT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Methodology
# ═════════════════════════════════════════════════════════════════════════════
def slide_06():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Methodology", "End-to-end pipeline for customer segmentation", 6)
    section_label(sl, "PIPELINE")
    slide_number(sl, 6)

    steps = [
        (NAVY,  "1\nData\nLoading",     "Load CC_GENERAL.csv\nInspect shape & dtypes\nCheck missing values"),
        (STEEL, "2\nExploratory\nData Analysis", "Distributions & KDE\nBoxplot outlier check\nCorrelation heatmap"),
        (TEAL,  "3\nMissing Value\nHandling",   "Identify nulls\nMedian imputation\nVerify 0 nulls remain"),
        (GREEN, "4\nFeature\nScaling",  "log₁p on skewed cols\nStandardScaler (z-score)\nSave scaler to disk"),
        (ORANGE,"5\nClustering\nComparison",    "K-Means (k=2–10)\nAgg., DBSCAN, GMM\nScore 3 metrics each"),
        (RGBColor(0x76,0x44,0x8A), "6\nFinal Model\nSelection",  "Pick best by silhouette\nInterpretability override\nRefit on full data"),
        (RED,   "7\nSegment\nInterpretation",   "Profile each cluster\nName the archetypes\nBusiness insights"),
    ]

    step_w = 1.68
    step_h = 1.65
    gap = 0.18
    y_top = 1.45
    for i, (color, label, desc) in enumerate(steps):
        x = 0.3 + i * (step_w + gap)
        rect(sl, x, y_top, step_w, step_h, fill=color, rounded=True)
        txt(sl, label, x + 0.08, y_top + 0.08, step_w - 0.12, step_h - 0.12,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 6:
            txt(sl, "→", x + step_w + 0.01, y_top + 0.62, 0.2, 0.4,
                size=14, bold=True, color=STEEL, align=PP_ALIGN.CENTER)

    # Description boxes below steps
    y_desc = 3.35
    for i, (color, _, desc) in enumerate(steps):
        x = 0.3 + i * (step_w + gap)
        rect(sl, x, y_desc, step_w, 2.95, fill=LIGHT_BG, line_color=color, line_w=0.8)
        dy = y_desc + 0.12
        for line in desc.split("\n"):
            txt(sl, f"  • {line}", x + 0.05, dy, step_w - 0.08, 0.5,
                size=9.5, color=TEXT)
            dy += 0.45

    # Tools footer
    rect(sl, 0.3, 6.45, 12.73, 0.75, fill=NAVY)
    txt(sl, "Tools:  Python 3  ·  pandas  ·  NumPy  ·  scikit-learn  ·  SciPy  ·  Matplotlib  ·  Seaborn  ·  Joblib  |  RANDOM_STATE = 42",
        0.5, 6.56, 12.33, 0.52, size=11, color=RGBColor(0xA9, 0xCC, 0xE3),
        align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Preprocessing
# ═════════════════════════════════════════════════════════════════════════════
def slide_07():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Preprocessing", "Cleaning and preparing the data for clustering", 7)
    section_label(sl, "DATA PREPARATION")
    slide_number(sl, 7)

    cards = [
        (NAVY, "Step 1 — Drop Identifier",
         ["CUST_ID is an arbitrary string ID\n(e.g., 'C10001') with no predictive signal",
          "Keeping it would confuse distance-based algorithms",
          "After dropping: 17 numeric features remain",
          "Shape: 8,950 rows × 17 columns"]),
        (STEEL, "Step 2 — Missing Value Imputation",
         ["MINIMUM_PAYMENTS: 313 nulls (3.5%)",
          "CREDIT_LIMIT: 1 null (<0.01%)",
          "Both columns are right-skewed →\nmedian is more robust than mean",
          "Verified: 0 nulls remain after imputation"]),
        (TEAL, "Step 3 — Outlier Analysis (No Removal)",
         ["Boxplots reveal extreme outliers in\nmonetary columns (cash-advance, purchases)",
          "We do NOT remove outlier rows:\n→ Tails often = most interesting customers",
          "High-spenders and heavy cash-advance\nusers are valuable segments themselves",
          "Handled via log₁p transform + scaling"]),
        (GREEN, "Step 4 — log₁p Transform + StandardScaler",
         ["log₁p applied to 10 skewed monetary cols:\n(BALANCE, PURCHASES, CASH_ADVANCE…)",
          "log₁p(0) = 0  →  zeros are preserved exactly",
          "StandardScaler: zero mean, unit variance\non all 17 features after log transform",
          "Fitted scaler saved → reusable for new data"]),
    ]

    xs = [0.3, 3.55, 6.8, 10.05]
    for i, (color, title, bullets) in enumerate(cards):
        x = xs[i]
        rect(sl, x, 1.35, 3.0, 5.65, fill=LIGHT_BG, line_color=color, line_w=1.2)
        rect(sl, x, 1.35, 3.0, 0.48, fill=color)
        txt(sl, title, x + 0.12, 1.37, 2.76, 0.43, size=11.5, bold=True, color=WHITE)
        by = 1.94
        for b in bullets:
            txt(sl, f"•  {b}", x + 0.14, by, 2.72, 0.8, size=10.5, color=TEXT)
            by += (0.72 if "\n" in b else 0.62)

    # Summary bar
    rect(sl, 0.3, 7.08, 12.73, 0.3, fill=NAVY)
    txt(sl, "Result:  8,950 × 17  scaled matrix  (cc_scaled.npy)  ·  Scaler saved to  models/scaler.joblib",
        0.5, 7.10, 12.33, 0.26, size=9.5, color=WHITE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Graphs After Preprocessing
# ═════════════════════════════════════════════════════════════════════════════
def slide_08():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Graphs After Preprocessing", "Visualising the cleaned, scaled, and transformed data", 8)
    section_label(sl, "POST-PROCESSING VISUALISATION")
    slide_number(sl, 8)

    txt(sl, "Feature Boxplots (Outlier Profile)", 0.3, 1.35, 7.5, 0.34,
        size=12, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/feature_boxplots.png", 0.3, 1.72, 7.8)

    txt(sl, "PCA Cumulative Explained Variance", 8.3, 1.35, 4.7, 0.34,
        size=12, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/pca_cumulative_variance.png", 8.3, 1.72, 4.7)

    rect(sl, 0.3, 6.5, 12.73, 0.78, fill=LIGHT_BG, line_color=STEEL, line_w=0.8)
    txt(sl, "Key Takeaways:", 0.5, 6.57, 2.5, 0.4, size=11, bold=True, color=STEEL)
    txt(sl, "Boxplots confirm heavy right tails — extreme high-value customers in every monetary column.  "
            "PCA shows ~7 components capture 90 % of variance.  "
            "We cluster on all 17 scaled features (not on PCA output) to preserve full information.",
        3.1, 6.57, 9.85, 0.56, size=10.5, color=TEXT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — What Is K-Means?
# ═════════════════════════════════════════════════════════════════════════════
def slide_09():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "What Is K-Means Clustering?", "A centroid-based partitional clustering algorithm", 9)
    section_label(sl, "ALGORITHM EXPLANATION")
    slide_number(sl, 9)

    # Concept text
    rect(sl, 0.3, 1.35, 7.8, 5.65, fill=LIGHT_BG, line_color=NAVY, line_w=1.0)
    txt(sl, "Core Concept", 0.5, 1.42, 7.4, 0.36, size=14, bold=True, color=NAVY)
    txt(sl,
        "K-Means partitions n data points into k non-overlapping clusters by "
        "iteratively minimising the total within-cluster sum of squared distances "
        "(inertia / WCSS) from each point to its assigned centroid.",
        0.5, 1.85, 7.4, 0.95, size=11.5, color=TEXT)

    steps = ["Randomly initialise k centroids in feature space",
             "Assign each point to the nearest centroid (Euclidean distance)",
             "Recompute each centroid as the mean of its assigned points",
             "Repeat steps 2–3 until centroids no longer move (convergence)",
             "Repeat with multiple random seeds; keep the best-inertia result"]
    sy = 2.9
    for i, s in enumerate(steps):
        rect(sl, 0.45, sy, 0.36, 0.36, fill=NAVY, rounded=True)
        txt(sl, str(i+1), 0.46, sy + 0.05, 0.34, 0.28, size=10.5,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, s, 0.88, sy + 0.04, 6.9, 0.35, size=11, color=TEXT)
        sy += 0.5

    txt(sl, "⚠  Assumes spherical, similarly-sized clusters.  "
            "Sensitive to outliers and random initialisation.  "
            "Results vary with k choice.",
        0.5, 5.7, 7.4, 0.6, size=10, color=ORANGE, italic=True)

    # Parameters panel
    rect(sl, 8.4, 1.35, 4.65, 5.65, fill=WHITE, line_color=STEEL, line_w=1.2)
    rect(sl, 8.4, 1.35, 4.65, 0.42, fill=STEEL)
    txt(sl, "Key Parameters", 8.6, 1.38, 4.3, 0.36, size=14, bold=True, color=WHITE)

    params = [
        ("k", "Number of clusters to form.\nWe sweep k = 2 … 10."),
        ("n_init", "How many times to run with\ndifferent random seeds (= 10)."),
        ("Inertia\n(WCSS)", "Sum of squared distances from\neach point to its centroid.\nLower = tighter clusters."),
        ("Silhouette\nScore", "Measures how much more tightly\na point fits its own cluster\nvs. the nearest other cluster."),
        ("random_state", "Fixed at 42 for\nreproducibility."),
    ]
    py = 1.9
    for pname, pdesc in params:
        rect(sl, 8.5, py, 4.45, 0.86 + (0.28 if "\n" in pname else 0),
             fill=LIGHT_BG, line_color=PALE, line_w=0.5)
        txt(sl, pname, 8.62, py + 0.07, 1.1, 0.7, size=10.5, bold=True, color=NAVY)
        txt(sl, pdesc, 9.75, py + 0.07, 2.95, 0.75, size=10, color=TEXT)
        py += (1.0 if "\n" in pname else 0.88)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — K-Means Output
# ═════════════════════════════════════════════════════════════════════════════
def slide_10():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "K-Means — Results & Model Selection", "Elbow method, silhouette analysis, and final k selection", 10)
    section_label(sl, "K-MEANS OUTPUT")
    slide_number(sl, 10)

    txt(sl, "Inertia (Elbow Curve)", 0.3, 1.35, 4.4, 0.34, size=12, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/kmeans_elbow.png", 0.3, 1.72, 4.4)

    txt(sl, "Silhouette Score vs k", 4.9, 1.35, 4.4, 0.34, size=12, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/kmeans_silhouette.png", 4.9, 1.72, 4.4)

    txt(sl, "PCA-2D Scatter  (k = 4)", 9.5, 1.35, 3.55, 0.34, size=12, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/kmeans_pca_scatter.png", 9.5, 1.72, 3.55)

    # Metrics table
    txt(sl, "K-Means Metrics per k  (best = k=4, selected for interpretability)", 0.3, 5.1, 9.0, 0.34,
        size=11, bold=True, color=NAVY)
    headers = ["k", "Silhouette ↑", "Davies-Bouldin ↓", "Calinski-Harabasz ↑", "Note"]
    rows = [
        ["2", "0.2515", "1.507", "3,064.9", "Auto peak — too coarse"],
        ["3", "0.2244", "1.681", "2,642.9", ""],
        ["4 ★", "0.2166", "1.688", "2,260.4", "Selected — inertia elbow + interpretability"],
        ["5", "0.2183", "1.601", "2,129.6", ""],
        ["9", "0.2204", "1.424", "1,845.8", "Metric best (excluding k=2) — too many segments"],
    ]
    pptx_table(sl, headers, rows, 0.3, 5.52, 9.0, 1.72,
               col_ratios=[0.6, 1.4, 1.5, 1.8, 2.7])

    rect(sl, 9.5, 5.1, 3.55, 2.14, fill=LIGHT_BG, line_color=NAVY, line_w=1.0)
    rect(sl, 9.5, 5.1, 3.55, 0.38, fill=NAVY)
    txt(sl, "Why k = 4?", 9.65, 5.13, 3.2, 0.32, size=12, bold=True, color=WHITE)
    for i, line in enumerate([
        "Silhouette peaks at k=2 — but 2 clusters\nmasks distinct customer types",
        "k=4 sits at the inertia elbow",
        "4 clusters → 4 interpretable archetypes:\nDormant, Revolvers, Transactors,\nInstalment Buyers",
    ]):
        txt(sl, f"•  {line}", 9.62, 5.58 + i * 0.55, 3.3, 0.55, size=10, color=TEXT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — What Is DBSCAN?
# ═════════════════════════════════════════════════════════════════════════════
def slide_11():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "What Is DBSCAN?", "Density-Based Spatial Clustering of Applications with Noise", 11)
    section_label(sl, "ALGORITHM EXPLANATION")
    slide_number(sl, 11)

    rect(sl, 0.3, 1.35, 7.8, 5.65, fill=LIGHT_BG, line_color=RED, line_w=1.0)
    txt(sl, "Core Concept", 0.5, 1.42, 7.4, 0.36, size=14, bold=True, color=RED)
    txt(sl,
        "DBSCAN discovers clusters as dense regions of points separated by low-density "
        "gaps. It does NOT require specifying k in advance. Instead, it defines clusters "
        "by two parameters (ε and min_samples) and automatically labels points in "
        "sparse regions as noise (label = –1).",
        0.5, 1.87, 7.4, 1.0, size=11.5, color=TEXT)

    concepts = [
        ("Core Point", "Has ≥ min_samples neighbours within radius ε"),
        ("Border Point", "Within ε of a core point, but fewer neighbours"),
        ("Noise Point", "Not reachable from any core point → label = –1"),
        ("Cluster", "All density-connected core + border points together"),
    ]
    cy = 3.0
    for cname, cdesc in concepts:
        rect(sl, 0.45, cy, 1.5, 0.44, fill=RED, rounded=True)
        txt(sl, cname, 0.5, cy + 0.07, 1.4, 0.32, size=9.5, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, cdesc, 2.1, cy + 0.07, 5.8, 0.35, size=11, color=TEXT)
        cy += 0.62

    txt(sl, "⚠  Struggles in high-dimensional, uniformly dense data (like financial features).  "
            "ε is hard to tune — too small = everything is noise, too large = one giant cluster.",
        0.5, 5.7, 7.4, 0.6, size=10, color=ORANGE, italic=True)

    # Parameters
    rect(sl, 8.4, 1.35, 4.65, 5.65, fill=WHITE, line_color=RED, line_w=1.2)
    rect(sl, 8.4, 1.35, 4.65, 0.42, fill=RED)
    txt(sl, "Key Parameters", 8.6, 1.38, 4.3, 0.36, size=14, bold=True, color=WHITE)
    params = [
        ("ε (eps)", "Maximum neighbourhood radius.\nChosen from k-distance plot elbow.\nWe swept: 1.55, 1.88, 2.21, 2.65"),
        ("min_samples", "Min. points to form a core point.\nSet to 5 throughout."),
        ("label = –1", "Points that are noise — not\nassigned to any cluster.\nProblematic for business use."),
        ("k-Distance\nPlot", "Sort the 5th-nearest-neighbour\ndistance for every point.\nElbow → good ε candidate."),
    ]
    py = 1.9
    for pname, pdesc in params:
        h = 1.05 if "\n" in pname else 0.95
        rect(sl, 8.5, py, 4.45, h, fill=LIGHT_BG, line_color=PALE, line_w=0.5)
        txt(sl, pname, 8.62, py + 0.07, 1.1, 0.85, size=10.5, bold=True, color=RED)
        txt(sl, pdesc, 9.75, py + 0.07, 2.95, h - 0.12, size=10, color=TEXT)
        py += h + 0.06


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — DBSCAN Output
# ═════════════════════════════════════════════════════════════════════════════
def slide_12():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "DBSCAN — Results", "k-distance elbow, parameter sweep, and PCA visualisation", 12)
    section_label(sl, "DBSCAN OUTPUT")
    slide_number(sl, 12)

    txt(sl, "k-Distance Plot (k=5)  →  choosing ε", 0.3, 1.35, 5.8, 0.34,
        size=12, bold=True, color=RED)
    img(sl, f"{PLOTS}/dbscan_k_distance.png", 0.3, 1.72, 5.8)

    txt(sl, "PCA-2D Scatter  (ε = 2.21,  grey = noise)", 6.35, 1.35, 6.68, 0.34,
        size=12, bold=True, color=RED)
    img(sl, f"{PLOTS}/dbscan_pca_scatter.png", 6.35, 1.72, 6.68)

    # Results table
    txt(sl, "ε Sweep Results  (min_samples = 5)", 0.3, 5.1, 6.5, 0.34,
        size=11, bold=True, color=RED)
    headers = ["ε", "Clusters Found", "Noise Points", "Silhouette ↑", "Davies-Bouldin ↓"]
    rows = [
        ["1.55", "20", "1,029  (11.5%)", "−0.211", "n/a"],
        ["1.88", "12", "419   (4.7%)",   "−0.123", "n/a"],
        ["2.21 ★", "2",  "153   (1.7%)",   "0.063",  "1.251"],
        ["2.65", "1",  "34    (0.4%)",   "n/a",    "n/a"],
    ]
    pptx_table(sl, headers, rows, 0.3, 5.52, 6.5, 1.72,
               col_ratios=[0.8, 1.4, 1.5, 1.2, 1.5])

    rect(sl, 7.0, 5.1, 6.05, 2.14, fill=RGBColor(0xFF, 0xEE, 0xEC),
         line_color=RED, line_w=1.0)
    rect(sl, 7.0, 5.1, 6.05, 0.38, fill=RED)
    txt(sl, "Honest Assessment", 7.15, 5.13, 5.7, 0.32, size=12, bold=True, color=WHITE)
    for i, line in enumerate([
        "Best silhouette 0.063 — far below K-Means (0.217)",
        "Credit card data has no geometric density structure",
        "~1.7 % of customers labelled 'noise' — unusable",
        "ε is very hard to tune in 17-D feature space",
        "→  DBSCAN not selected as final model",
    ]):
        c = RED if line.startswith("→") else TEXT
        txt(sl, f"  {'✗' if not line.startswith('→') else ''}  {line}",
            7.1, 5.57 + i * 0.33, 5.8, 0.35, size=10, color=c)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — What Is GMM?
# ═════════════════════════════════════════════════════════════════════════════
def slide_13():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "What Is Gaussian Mixture Model (GMM)?", "Probabilistic soft-assignment clustering", 13)
    section_label(sl, "ALGORITHM EXPLANATION")
    slide_number(sl, 13)

    rect(sl, 0.3, 1.35, 7.8, 5.65, fill=LIGHT_BG, line_color=TEAL, line_w=1.0)
    txt(sl, "Core Concept", 0.5, 1.42, 7.4, 0.36, size=14, bold=True, color=TEAL)
    txt(sl,
        "GMM models the data as a mixture of k multivariate Gaussian distributions "
        "(ellipsoidal clusters). Unlike K-Means which makes hard assignments, GMM "
        "gives each point a probability of belonging to each component — this is "
        "called soft clustering. The EM algorithm is used to fit the model.",
        0.5, 1.87, 7.4, 1.1, size=11.5, color=TEXT)

    concepts = [
        ("Component", "One Gaussian (ellipsoid) in the mixture — each has its own mean, covariance, and weight"),
        ("Soft\nAssignment", "Each point gets a probability vector across all k components.\nFinal label = component with highest probability"),
        ("EM Algorithm", "Expectation-Maximisation iteratively estimates component parameters until convergence"),
        ("covariance\ntype='full'", "Each component has its own full covariance matrix — the most flexible but expensive option"),
    ]
    cy = 3.1
    for cname, cdesc in concepts:
        h = 0.72 if "\n" in cname else 0.6
        rect(sl, 0.45, cy, 1.5, h, fill=TEAL, rounded=True)
        txt(sl, cname, 0.5, cy + 0.1, 1.4, h - 0.12, size=9.5, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, cdesc, 2.1, cy + 0.07, 5.8, h - 0.1, size=11, color=TEXT)
        cy += h + 0.14

    txt(sl, "⚠  More expressive than K-Means but uses far more parameters.  "
            "BIC and AIC penalise complexity and help select the right number of components.",
        0.5, 5.72, 7.4, 0.55, size=10, color=ORANGE, italic=True)

    rect(sl, 8.4, 1.35, 4.65, 5.65, fill=WHITE, line_color=TEAL, line_w=1.2)
    rect(sl, 8.4, 1.35, 4.65, 0.42, fill=TEAL)
    txt(sl, "Key Parameters", 8.6, 1.38, 4.3, 0.36, size=14, bold=True, color=WHITE)
    params = [
        ("n_components", "Number of Gaussian components.\nWe sweep n = 2 … 8."),
        ("covariance\ntype", "Shape of each cluster ellipsoid.\n'full' = each has own covariance."),
        ("BIC", "Bayesian Information Criterion.\nLower BIC = better fit with\ncomplexity penalty."),
        ("AIC", "Akaike Information Criterion.\nSimilar to BIC but weaker\ncomplexity penalty."),
        ("max_iter", "Max EM iterations = 200.\nrandom_state = 42."),
    ]
    py = 1.9
    for pname, pdesc in params:
        h = 1.0 if "\n" in pname else 0.88
        rect(sl, 8.5, py, 4.45, h, fill=LIGHT_BG, line_color=PALE, line_w=0.5)
        txt(sl, pname, 8.62, py + 0.07, 1.2, h - 0.1, size=10.5, bold=True, color=TEAL)
        txt(sl, pdesc, 9.85, py + 0.07, 2.85, h - 0.12, size=10, color=TEXT)
        py += h + 0.06


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — GMM Output
# ═════════════════════════════════════════════════════════════════════════════
def slide_14():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "GMM — Results", "BIC/AIC curves, metrics, and PCA visualisation", 14)
    section_label(sl, "GMM OUTPUT")
    slide_number(sl, 14)

    txt(sl, "BIC and AIC vs n_components", 0.3, 1.35, 6.5, 0.34,
        size=12, bold=True, color=TEAL)
    img(sl, f"{PLOTS}/gmm_bic_aic.png", 0.3, 1.72, 6.5)

    txt(sl, "PCA-2D Scatter  (n = 2, best silhouette)", 7.0, 1.35, 6.05, 0.34,
        size=12, bold=True, color=TEAL)
    img(sl, f"{PLOTS}/gmm_pca_scatter.png", 7.0, 1.72, 6.05)

    txt(sl, "GMM Metrics per n_components", 0.3, 5.1, 6.5, 0.34,
        size=11, bold=True, color=TEAL)
    headers = ["n", "Silhouette ↑", "Davies-Bouldin ↓", "Calinski-Harabasz ↑", "BIC ↓"]
    rows = [
        ["2 ★", "0.2105", "1.826", "2,382", "66,352"],
        ["3", "0.1666", "2.168", "2,209", "−71,190"],
        ["4", "0.1537", "1.954", "1,956", "−114,683"],
        ["6", "0.1703", "1.745", "1,888", "−255,534"],
        ["8", "0.1449", "1.901", "1,642", "−322,300"],
    ]
    pptx_table(sl, headers, rows, 0.3, 5.52, 6.5, 1.72,
               col_ratios=[0.6, 1.2, 1.5, 1.8, 1.0])

    rect(sl, 7.0, 5.1, 6.05, 2.14, fill=LIGHT_BG, line_color=TEAL, line_w=1.0)
    rect(sl, 7.0, 5.1, 6.05, 0.38, fill=TEAL)
    txt(sl, "GMM vs. K-Means", 7.15, 5.13, 5.7, 0.32, size=12, bold=True, color=WHITE)
    for i, line in enumerate([
        "Best silhouette 0.2105 at n=2 (vs K-Means 0.2166)",
        "BIC decreases monotonically — no clear optimum",
        "Soft assignments add no interpretability benefit here",
        "K-Means centroids are more readable as profiles",
        "→  GMM competitive but K-Means wins on interpretability",
    ]):
        c = TEAL if line.startswith("→") else TEXT
        txt(sl, f"  {'•' if not line.startswith('→') else '★'}  {line}",
            7.1, 5.56 + i * 0.33, 5.8, 0.35, size=10, color=c)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — What Is Agglomerative Clustering?
# ═════════════════════════════════════════════════════════════════════════════
def slide_15():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "What Is Agglomerative Clustering?", "Hierarchical bottom-up clustering with Ward linkage", 15)
    section_label(sl, "ALGORITHM EXPLANATION")
    slide_number(sl, 15)

    rect(sl, 0.3, 1.35, 7.8, 5.65, fill=LIGHT_BG, line_color=GREEN, line_w=1.0)
    txt(sl, "Core Concept", 0.5, 1.42, 7.4, 0.36, size=14, bold=True, color=GREEN)
    txt(sl,
        "Agglomerative Clustering builds a hierarchy of clusters from the bottom up. "
        "It starts with every point as its own cluster, then iteratively merges the "
        "two closest clusters until only k clusters remain. The merge history can be "
        "visualised as a dendrogram, which helps choose k visually.",
        0.5, 1.87, 7.4, 1.1, size=11.5, color=TEXT)

    steps = [
        "Start: each of the 8,950 customers is its own cluster",
        "Find the pair of clusters with the smallest merge distance",
        "Merge them into one cluster (Ward: minimise total within-cluster variance)",
        "Repeat until exactly k clusters remain",
        "The dendrogram records every merge — large jumps = natural cut points",
    ]
    sy = 3.1
    for i, s in enumerate(steps):
        rect(sl, 0.45, sy, 0.36, 0.36, fill=GREEN, rounded=True)
        txt(sl, str(i+1), 0.46, sy + 0.05, 0.34, 0.28, size=10.5,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, s, 0.88, sy + 0.04, 6.9, 0.35, size=11, color=TEXT)
        sy += 0.5

    txt(sl, "⚠  Does not scale well to very large datasets (O(n²) memory).  "
            "No predict() method — labels come from labels_ attribute only.",
        0.5, 5.72, 7.4, 0.55, size=10, color=ORANGE, italic=True)

    rect(sl, 8.4, 1.35, 4.65, 5.65, fill=WHITE, line_color=GREEN, line_w=1.2)
    rect(sl, 8.4, 1.35, 4.65, 0.42, fill=GREEN)
    txt(sl, "Key Parameters", 8.6, 1.38, 4.3, 0.36, size=14, bold=True, color=WHITE)
    params = [
        ("n_clusters", "Number of clusters to cut the\nhierarchy at.\nWe sweep n = 2 … 8."),
        ("linkage\n= 'ward'", "How to measure distance between\nclusters. Ward minimises total\nwithin-cluster variance."),
        ("Dendrogram", "Tree diagram of merge history.\nLarge vertical jumps indicate\nnatural cluster boundaries."),
        ("Affinity", "Distance metric (default:\nEuclidean). Ward linkage\nrequires Euclidean."),
    ]
    py = 1.9
    for pname, pdesc in params:
        h = 1.05 if "\n" in pname else 0.92
        rect(sl, 8.5, py, 4.45, h, fill=LIGHT_BG, line_color=PALE, line_w=0.5)
        txt(sl, pname, 8.62, py + 0.07, 1.2, h - 0.1, size=10.5, bold=True, color=GREEN)
        txt(sl, pdesc, 9.85, py + 0.07, 2.85, h - 0.12, size=10, color=TEXT)
        py += h + 0.06


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Agglomerative Output
# ═════════════════════════════════════════════════════════════════════════════
def slide_16():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Agglomerative Clustering — Results", "Ward-linkage dendrogram and PCA visualisation", 16)
    section_label(sl, "AGGLOMERATIVE OUTPUT")
    slide_number(sl, 16)

    txt(sl, "Dendrogram  (500-row sample, Ward linkage)", 0.3, 1.35, 6.9, 0.34,
        size=12, bold=True, color=GREEN)
    img(sl, f"{PLOTS}/agglomerative_dendrogram.png", 0.3, 1.72, 6.9)

    txt(sl, "PCA-2D Scatter  (n = 2, best silhouette)", 7.4, 1.35, 5.65, 0.34,
        size=12, bold=True, color=GREEN)
    img(sl, f"{PLOTS}/agglomerative_pca_scatter.png", 7.4, 1.72, 5.65)

    txt(sl, "Agglomerative Metrics per n", 0.3, 5.1, 6.9, 0.34,
        size=11, bold=True, color=GREEN)
    headers = ["n", "Silhouette ↑", "Davies-Bouldin ↓", "Calinski-Harabasz ↑"]
    rows = [
        ["2 ★", "0.2154", "1.522", "2,251"],
        ["3", "0.1745", "1.896", "1,978"],
        ["4", "0.1658", "1.733", "1,734"],
        ["7", "0.1535", "1.744", "1,521"],
        ["8", "0.1632", "1.678", "1,498"],
    ]
    pptx_table(sl, headers, rows, 0.3, 5.52, 6.9, 1.72,
               col_ratios=[0.8, 1.4, 1.5, 1.8])

    rect(sl, 7.4, 5.1, 5.65, 2.14, fill=LIGHT_BG, line_color=GREEN, line_w=1.0)
    rect(sl, 7.4, 5.1, 5.65, 0.38, fill=GREEN)
    txt(sl, "Assessment", 7.55, 5.13, 5.3, 0.32, size=12, bold=True, color=WHITE)
    for i, line in enumerate([
        "Best silhouette 0.2154 at n=2 — close to K-Means",
        "Ward linkage gives clean, balanced dendrogram",
        "Best DB (1.522) — clusters are well-separated",
        "But n=2 is too coarse; n=4 silhouette drops to 0.166",
        "→  K-Means k=4 preferred for interpretable segments",
    ]):
        c = GREEN if line.startswith("→") else TEXT
        txt(sl, f"  {'•' if not line.startswith('→') else '★'}  {line}",
            7.5, 5.56 + i * 0.33, 5.4, 0.35, size=10, color=c)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Comparison and Final Algorithm Selection
# ═════════════════════════════════════════════════════════════════════════════
def slide_17():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Algorithm Comparison & Final Selection", "Choosing the winner using metrics + interpretability", 17)
    section_label(sl, "MODEL SELECTION")
    slide_number(sl, 17)

    txt(sl, "Best Configuration per Algorithm", 0.3, 1.35, 7.5, 0.34,
        size=12, bold=True, color=NAVY)
    headers = ["Algorithm", "Config", "Silhouette ↑", "Davies-Bouldin ↓",
               "Calinski-Harabasz ↑", "Noise", "BIC"]
    rows = [
        ["K-Means  ★", "k = 4", "0.2166", "1.688", "2,260", "—", "—"],
        ["Agglomerative", "n=2, Ward", "0.2154", "1.522", "2,251", "—", "—"],
        ["GMM", "n=2, full cov.", "0.2105", "1.826", "2,382", "—", "66,352"],
        ["DBSCAN", "ε=2.21, min=5", "0.0634", "1.251", "3", "153", "—"],
    ]
    pptx_table(sl, headers, rows, 0.3, 1.78, 9.0, 1.92,
               col_ratios=[1.7, 1.5, 1.2, 1.5, 1.7, 0.7, 0.9])

    img(sl, f"{PLOTS}/algorithm_comparison_silhouette.png", 9.5, 1.35, 3.55)

    # Decision reasoning
    reasons = [
        (NAVY,   "K-Means Wins on Metrics",
         "Highest silhouette (0.2166) + Highest Calinski-Harabasz (2,260).\nCentroids are directly readable as average customer behaviour."),
        (RED,    "DBSCAN Eliminated",
         "Silhouette 0.063 — far below all others.\nLabels 153 customers as noise — unusable for business."),
        (TEAL,   "GMM Competitive but Rejected",
         "Silhouette slightly lower than K-Means.\nNo interpretability advantage from soft assignments here."),
        (GREEN,  "Agglomerative — Good DB, Wrong k",
         "Best Davies-Bouldin (1.522) at n=2, but n=4 drops to 0.166.\nNo clear elbow to justify 4 segments."),
    ]
    rx = [0.3, 3.55, 6.8, 9.5]
    for i, (color, title, desc) in enumerate(reasons):
        rect(sl, rx[i], 3.9, 3.0, 2.95, fill=LIGHT_BG, line_color=color, line_w=1.0)
        rect(sl, rx[i], 3.9, 3.0, 0.42, fill=color)
        txt(sl, title, rx[i] + 0.12, 3.93, 2.76, 0.36, size=11, bold=True, color=WHITE)
        txt(sl, desc, rx[i] + 0.14, 4.42, 2.72, 1.8, size=10, color=TEXT)

    rect(sl, 0.3, 7.05, 12.73, 0.32, fill=NAVY)
    txt(sl, "★  Final model:  K-Means  k = 4   (silhouette=0.2166,  DB=1.688,  CH=2,260)  —  interpretability override from auto-peak at k=2",
        0.5, 7.07, 12.33, 0.26, size=10, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Modeling
# ═════════════════════════════════════════════════════════════════════════════
def slide_18():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Modeling — Fitting the Final Model", "Refitting K-Means k=4 cleanly on the full dataset", 18)
    section_label(sl, "FINAL MODEL")
    slide_number(sl, 18)

    steps = [
        (NAVY,  "1. Load Best Choice",
         "Read outputs/results/best_model_choice.json\nConfirm: algorithm=KMeans, params='k=4'"),
        (STEEL, "2. Instantiate Model",
         "KMeans(n_clusters=4,\n  random_state=42,\n  n_init=10)"),
        (TEAL,  "3. Fit on Full Data",
         "model.fit(X_scaled)\nX_scaled = cc_scaled.npy\n(8,950 × 17 matrix)"),
        (GREEN, "4. Get Labels",
         "labels = model.predict(X)\nReturns cluster ID (0,1,2,3)\nfor each of 8,950 customers"),
        (ORANGE,"5. Compute Final\nMetrics",
         "Silhouette:   0.2166\nDavies-Bouldin: 1.688\nCalinski-Harabasz: 2,260"),
        (RGBColor(0x76,0x44,0x8A), "6. Persist Results",
         "models/final_model.joblib\nmodels/cluster_labels.npy\noutputs/results/final_metrics.json"),
    ]

    xs = [0.3, 2.55, 4.8, 7.05, 9.3, 11.0]
    ws = [2.0, 2.0, 2.0, 2.0, 2.0, 2.05]
    for i, ((color, title, desc), x, w) in enumerate(zip(steps, xs, ws)):
        rect(sl, x, 1.35, w, 4.8, fill=LIGHT_BG, line_color=color, line_w=1.1)
        rect(sl, x, 1.35, w, 0.5, fill=color)
        txt(sl, title, x + 0.1, 1.37, w - 0.15, 0.46,
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x + 0.08, 2.0, w - 0.14, 3.0, fill=WHITE, line_color=color, line_w=0.5)
        txt(sl, desc, x + 0.15, 2.1, w - 0.25, 2.85, size=10, color=TEXT,
            font="Courier New")
        if i < 5:
            txt(sl, "→", x + w + 0.02, 2.65, 0.2, 0.4, size=14,
                bold=True, color=color, align=PP_ALIGN.CENTER)

    # Outcome banner
    rect(sl, 0.3, 6.35, 12.73, 1.0, fill=NAVY)
    txt(sl, "Outcome", 0.5, 6.4, 2.0, 0.35, size=12, bold=True,
        color=RGBColor(0xA9, 0xCC, 0xE3))
    txt(sl, "Each of the 8,950 customers receives a cluster label (0 – 3).  "
            "Labels attached to cc_clean.csv (original un-scaled values) for interpretable profiling.  "
            "Model and scaler persisted for scoring new customers.",
        2.6, 6.4, 10.3, 0.82, size=11.5, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Final Results
# ═════════════════════════════════════════════════════════════════════════════
def slide_19():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    slide_header(sl, "Final Results — Customer Segments", "Four interpretable archetypes from K-Means k=4", 19)
    section_label(sl, "RESULTS")
    slide_number(sl, 19)

    # Three plots
    txt(sl, "Cluster Sizes", 0.3, 1.35, 3.8, 0.3, size=11, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/final_cluster_sizes.png", 0.3, 1.68, 3.8)

    txt(sl, "PCA-2D Cluster Visualization", 4.3, 1.35, 4.3, 0.3, size=11, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/final_clusters_pca.png", 4.3, 1.68, 4.3)

    txt(sl, "Cluster Profile Heatmap", 8.8, 1.35, 4.25, 0.3, size=11, bold=True, color=NAVY)
    img(sl, f"{PLOTS}/cluster_profile_heatmap.png", 8.8, 1.68, 4.25)

    # Segment cards
    segments = [
        (CC[0], "Cluster 0", "Dormant / Light Users",
         "~48%  ·  ~4,300 customers",
         ["Avg Balance: $214", "Avg Purchases: $351",
          "Purchase Freq: 0.25", "→ Rarely use the card"]),
        (CC[1], "Cluster 1", "Cash-Advance Revolvers",
         "~13%  ·  ~1,200 customers",
         ["Cash Advance: $2,220", "Full-Pay Rate: 0.03",
          "Balance: $2,461", "→ Card used as a loan"]),
        (CC[2], "Cluster 2", "Premium Transactors",
         "~19%  ·  ~1,730 customers",
         ["Purchases: $2,604", "Credit Limit: $6,393",
          "~35 TRX/month", "→ High-value, high-spend"]),
        (CC[3], "Cluster 3", "Instalment Buyers",
         "~19%  ·  ~1,720 customers",
         ["Instalment Spend: $631", "Instalment Freq: 0.76",
          "Full-Pay Rate: 0.29", "→ Disciplined spenders"]),
    ]
    xs = [0.3, 3.55, 6.8, 10.05]
    for i, (color, cid, cname, size, bullets) in enumerate(segments):
        x = xs[i]
        rect(sl, x, 5.2, 3.0, 2.1, fill=LIGHT_BG, line_color=color, line_w=1.2)
        rect(sl, x, 5.2, 3.0, 0.6, fill=color)
        txt(sl, cid, x + 0.12, 5.22, 1.0, 0.28, size=9.5, bold=True, color=WHITE)
        txt(sl, cname, x + 0.12, 5.44, 2.76, 0.32, size=11, bold=True, color=WHITE)
        txt(sl, size, x + 0.12, 5.78, 2.76, 0.24, size=9, color=LGRAY, italic=True)
        by = 6.08
        for b in bullets:
            c = color if b.startswith("→") else TEXT
            txt(sl, f"  {b}", x + 0.08, by, 2.82, 0.3, size=9.5, color=c)
            by += 0.3


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Conclusion
# ═════════════════════════════════════════════════════════════════════════════
def slide_20():
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    rect(sl, 0, 0, 0.12, 7.5, fill=STEEL)
    rect(sl, 0, 0, 13.33, 0.06, fill=STEEL)
    rect(sl, 0, 7.44, 13.33, 0.06, fill=STEEL)

    txt(sl, "Conclusion", 0.4, 0.25, 8.0, 0.7,
        size=32, bold=True, color=WHITE)
    rect(sl, 0.4, 0.95, 8.0, 0.04, fill=STEEL)

    # Summary column
    rect(sl, 0.4, 1.12, 6.0, 5.6, fill=RGBColor(0x1E, 0x44, 0x6A))
    txt(sl, "What We Achieved", 0.6, 1.2, 5.6, 0.38, size=14, bold=True,
        color=RGBColor(0xA9, 0xCC, 0xE3))
    summary = [
        "Applied 4 clustering algorithms to 8,950 credit card customers",
        "Evaluated each with Silhouette, Davies-Bouldin, and Calinski-Harabasz",
        "Selected K-Means k=4 — best metrics + highest interpretability",
        "Discovered 4 distinct behavioral archetypes in the customer base",
        "Mapped each segment to specific banking products and strategies",
        "Persisted fitted model and scaler for deployment on new customers",
    ]
    sy = 1.68
    for s in summary:
        rect(sl, 0.52, sy, 0.28, 0.28, fill=STEEL, rounded=True)
        txt(sl, "✓", 0.53, sy + 0.04, 0.26, 0.22, size=9, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, s, 0.88, sy + 0.02, 5.3, 0.46, size=11, color=WHITE)
        sy += 0.52

    # Business value column
    rect(sl, 6.65, 1.12, 6.35, 2.7, fill=RGBColor(0x1E, 0x44, 0x6A))
    txt(sl, "Business Value", 6.85, 1.2, 5.9, 0.38, size=14, bold=True,
        color=RGBColor(0xA9, 0xCC, 0xE3))
    bv = [("C0 Dormant",    "Re-engagement campaigns, fee waivers"),
          ("C1 Revolvers",  "Balance-transfer, risk monitoring"),
          ("C2 Transactors","Premium rewards, credit limit upgrades"),
          ("C3 Instalment", "BNPL cross-sell, low-rate plans")]
    colors = CC
    by = 1.68
    for (lbl, action), c in zip(bv, colors):
        rect(sl, 6.75, by, 1.4, 0.46, fill=c, rounded=True)
        txt(sl, lbl, 6.8, by + 0.08, 1.3, 0.3, size=9.5, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, action, 8.22, by + 0.08, 4.6, 0.36, size=10.5, color=WHITE)
        by += 0.56

    # Future work
    rect(sl, 6.65, 4.02, 6.35, 2.8, fill=RGBColor(0x1E, 0x44, 0x6A))
    txt(sl, "Future Improvements", 6.85, 4.1, 5.9, 0.38, size=14, bold=True,
        color=RGBColor(0xA9, 0xCC, 0xE3))
    future = [
        "Apply log₁p transform as part of the pipeline (already done in this project)",
        "Experiment with outlier capping (IQR method) before scaling",
        "Try more advanced metrics: Gap Statistic, Hopkins test for clusterability",
        "Enrich with demographic data (age, income, region) for richer segments",
        "Deploy model as REST API — score new customers at onboarding",
        "Re-run segmentation quarterly to detect segment migration over time",
    ]
    fy = 4.55
    for f in future:
        txt(sl, f"  →  {f}", 6.75, fy, 6.1, 0.38, size=10, color=RGBColor(0xCC, 0xE0, 0xF0))
        fy += 0.38

    # Bottom credits bar
    rect(sl, 0, 6.92, 13.33, 0.58, fill=RGBColor(0x0F, 0x24, 0x38))
    txt(sl, "University of Delaware  ·  Introduction to Data Mining  ·  May 2026  "
            "·  K-Means k=4  ·  Silhouette = 0.2166  ·  4 Customer Segments",
        0.4, 6.99, 12.53, 0.4, size=10,
        color=RGBColor(0x5C, 0x85, 0xAD), align=PP_ALIGN.CENTER)


# ── Build all slides ──────────────────────────────────────────────────────────
slide_01(); slide_02(); slide_03(); slide_04(); slide_05()
slide_06(); slide_07(); slide_08(); slide_09(); slide_10()
slide_11(); slide_12(); slide_13(); slide_14(); slide_15()
slide_16(); slide_17(); slide_18(); slide_19(); slide_20()

prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
